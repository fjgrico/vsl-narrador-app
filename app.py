# ✅ app.py COMPLETO Y FUNCIONAL PARA STREAMLIT CLOUD

import streamlit as st
import openai
import os
import requests
from gtts import gTTS
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches

# ✅ Cargar claves desde secrets
try:
    openai.api_key = st.secrets["OPENAI_API_KEY"]
    openai_project_id = st.secrets["OPENAI_PROJECT_ID"]
    eleven_api_key = st.secrets["ELEVEN_API_KEY"]
except KeyError:
    st.error("❌ No se encontraron las claves API en los secrets. Revisa la configuración.")
    st.stop()

# ✅ Función para generar guion
@st.cache_data(show_spinner=True)
def generar_guion(producto, publico, problema, beneficios, objeciones, garantia, precio, llamada):
    prompt = f"""
Eres un experto en copywriting y VSL. Crea un guion de máximo 5 minutos, directo, sencillo, profesional, sin tecnicismos.

1. Hook inicial
2. Historia personal (breve)
3. Problema (con emoción)
4. Presentación de la solución
5. Beneficios ({beneficios})
6. Objeciones típicas y cómo superarlas ({objeciones})
7. Garantía ofrecida ({garantia})
8. Precio: {precio}
9. Llamada a la acción: {llamada}

Producto: {producto}
Público objetivo: {publico}
Problema que resuelve: {problema}
"""
    client = openai.OpenAI(api_key=openai.api_key)
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    return response.choices[0].message.content

# ✅ Función para generar audio con gTTS
def generar_audio(texto):
    tts = gTTS(text=texto, lang='es')
    filename = "narracion.mp3"
    tts.save(filename)
    return filename

# ✅ Función para generar PDF
def generar_pdf(texto):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for line in texto.split("\n"):
        pdf.multi_cell(0, 10, line)
    pdf.output("guion.pdf")
    return "guion.pdf"

# ✅ Función para generar PowerPoint
def generar_ppt(texto):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    for i, seccion in enumerate(texto.split("\n\n")):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Sección {i+1}"
        slide.placeholders[1].text = seccion.strip()
    prs.save("guion.pptx")
    return "guion.pptx"

# ✅ Interfaz principal
st.title("🎤 Generador VSL con Narración por Voz IA")
st.markdown("Completa los datos y genera tu guion con audio, PDF y PowerPoint.")

with st.form("formulario"):
    producto = st.text_input("Producto o servicio")
    publico = st.text_input("Público objetivo")
    problema = st.text_area("¿Qué problema resuelve?")
    beneficios = st.text_area("Beneficios (separados por comas)")
    objeciones = st.text_area("Objeciones típicas")
    garantia = st.text_input("¿Qué garantía ofreces?")
    precio = st.text_input("Precio o forma de pago")
    llamada = st.text_input("Llamada a la acción")
    submitted = st.form_submit_button("🧠 Generar Guion")

# ✅ Mostrar mensaje al cargar la app si no se ha enviado el formulario
if not submitted:
    st.info("👈 Completa los datos del formulario y pulsa 'Generar Guion' para comenzar.")

if submitted:
    with st.spinner("Generando guion..."):
        guion = generar_guion(producto, publico, problema, beneficios, objeciones, garantia, precio, llamada)
        st.subheader("📝 Guion Generado")
        st.write(guion)

        pdf_path = generar_pdf(guion)
        ppt_path = generar_ppt(guion)
        audio_path = generar_audio(guion)

        st.download_button("📄 Descargar TXT", guion, file_name="guion.txt")
        st.download_button("📄 Descargar PDF", open(pdf_path, "rb"), file_name="guion.pdf")
        st.download_button("📊 Descargar PowerPoint", open(ppt_path, "rb"), file_name="guion.pptx")
        st.audio(open(audio_path, "rb"), format="audio/mp3")